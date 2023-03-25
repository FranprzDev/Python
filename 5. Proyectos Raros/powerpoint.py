# Importar librerías necesarias
from pptx import Presentation
from pptx.util import Inches

# Crear presentación
prs = Presentation()
Container = collections.abc.Container

# Agregar diapositiva de título
titulo_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(titulo_slide_layout)
titulo = slide.shapes.title
subtitulo = slide.placeholders[1]
titulo.text = "Los beneficios de la naturaleza"
subtitulo.text = "Por ChatGPT"

# Agregar diapositiva de introducción
introduccion_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(introduccion_slide_layout)
titulo = slide.shapes.title
subtitulo = slide.placeholders[1]
titulo.text = "La naturaleza es nuestra aliada"
subtitulo.text = "Breve introducción sobre la importancia de la naturaleza para la salud y el bienestar humano"
imagen_path = "ruta/de/la/imagen.jpg"
imagen = slide.shapes.add_picture(imagen_path, Inches(1), Inches(1), height=Inches(4), width=Inches(6))

# Agregar diapositiva de beneficios para la salud física
fisica_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(fisica_slide_layout)
titulo = slide.shapes.title
subtitulo = slide.placeholders[1]
titulo.text = "La naturaleza como fuente de bienestar físico"
subtitulo.text = "Descripción de cómo la naturaleza puede ayudar a mejorar la salud física, como la reducción del estrés, el aumento de la actividad física y la mejora del sistema inmunológico"
imagen_path = "ruta/de/la/imagen.jpg"
imagen = slide.shapes.add_picture(imagen_path, Inches(1), Inches(1), height=Inches(4), width=Inches(6))

# Agregar diapositiva de beneficios para la salud mental
mental_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(mental_slide_layout)
titulo = slide.shapes.title
subtitulo = slide.placeholders[1]
titulo.text = "La naturaleza como fuente de bienestar emocional"
subtitulo.text = "Descripción de cómo la naturaleza puede ayudar a mejorar la salud mental, como la reducción de la ansiedad, la depresión y la mejora del estado de ánimo"
imagen_path = "ruta/de/la/imagen.jpg"
imagen = slide.shapes.add_picture(imagen_path, Inches(1), Inches(1), height=Inches(4), width=Inches(6))

# Agregar diapositiva de beneficios para el medio ambiente
ambiente_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(ambiente_slide_layout)
titulo = slide.shapes.title
subtitulo = slide.placeholders[1]
titulo.text = "La naturaleza como fuente de sustentabilidad"
subtitulo.text = "Descripción de cómo la naturaleza puede ayudar a preservar el medio ambiente y promover la sustentabilidad, como la reducción de la contaminación y la protección de la biodiversidad"
imagen_path = "ruta/de/la/imagen.jpg"
imagen = slide.shapes.add_picture(imagen_path, Inches(1), Inches(1), height=Inches(4), width=Inches(6))

prs.save("presentacion_naturaleza.pptx")